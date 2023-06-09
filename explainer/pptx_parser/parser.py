from pptx import Presentation


def parse_presentation(file_path: str) -> list:
    """
    Parse a PowerPoint presentation and extract the text from each slide.

    This function opens a PowerPoint file specified by `file_path`, and for each slide in the presentation, it extracts
    all the text content and appends it to a list. The lists of slide text are then added to a master list, which is
    returned.

    Args:
        file_path (str): The path to the PowerPoint file to parse.

    Returns:
        list: A list of lists, where each sublist contains the text content of a single slide in the presentation.
    """
    presentation = Presentation(file_path)
    slides = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
        slides.append(slide_text)
    return slides


if __name__ == '__main__':

    parsed_data = parse_presentation("asyncio-intro.pptx")
    for i, _slide in enumerate(parsed_data):
        print(f"Slide {i + 1}:")
        print(_slide)
