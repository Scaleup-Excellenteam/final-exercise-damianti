from pptx import Presentation


def parse_presentation(file_path: str) -> list:
    presentation = Presentation(file_path)
    slides = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
        slides.append(slide_text)
    return slides


if __name__ == '__main__':

    parsed_data = parse_presentation("asyncio-intro.pptx")
    for i, _slide in enumerate(parsed_data):
        print(f"Slide {i + 1}:")
        print(_slide)
